from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor  # Import RGBColor for color assignment

# Initialize a new presentation
prs = Presentation()

# -------- 1. Cover Slide --------
slide_cover = prs.slides.add_slide(prs.slide_layouts[6])
left, top, width, height = Inches(0), Inches(1), prs.slide_width, Inches(1.2)
# Name
name_shape = slide_cover.shapes.add_textbox(left, Inches(0.6), width, height)
name_tf = name_shape.text_frame
name_p = name_tf.add_paragraph()
name_p.text = "이상훈 (SangHoon Lee)"
name_p.font.size = Pt(40)
name_p.font.bold = True
name_p.alignment = PP_ALIGN.CENTER
# Job Title
role_shape = slide_cover.shapes.add_textbox(left, Inches(2.0), width, height)
role_tf = role_shape.text_frame
role_p = role_tf.add_paragraph()
role_p.text = "Full Stack AI Developer  |  Solution Architect"
role_p.font.size = Pt(24)
role_p.font.italic = True
role_p.alignment = PP_ALIGN.CENTER
# Brief Tagline
tagline_shape = slide_cover.shapes.add_textbox(left, Inches(2.8), width, height)
tagline_tf = tagline_shape.text_frame
tagline_p = tagline_tf.add_paragraph()
tagline_p.text = "AI 자동화 기반 IT 서비스 설계 및 데이터 기반 플랫폼 구축 전문가"
tagline_p.font.size = Pt(18)
tagline_p.alignment = PP_ALIGN.CENTER

# Optional: Contact Info area
contact_shape = slide_cover.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(10), Inches(0.5))
contact_tf = contact_shape.text_frame
contact_tf.word_wrap = True
contact = contact_tf.add_paragraph()
contact.text = "📧 Email: itsme.sanghoon@gmail.com   |   🌐 github.com/shlee-ai"
contact.font.size = Pt(14)
contact.alignment = PP_ALIGN.CENTER

# -------- 2. Professional Summary Slide --------
slide_summary = prs.slides.add_slide(prs.slide_layouts[5])
slide_summary.shapes.title.text = "Professional Summary"

# -- Avoid KeyError by not using shapes.placeholders[1]; use add_textbox instead
summary_box = slide_summary.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(1.6))
summary_tf = summary_box.text_frame
summary_tf.word_wrap = True
summary_text = (
    "• 생성형 AI 및 데이터 분석 중심 풀스택 개발자로서 10년 이상 실전 경험\n"
    "• Python, Django, React, Streamlit 등 최신 기술을 통한 E2E 플랫폼 구축\n"
    "• AI 자동화, 데이터 파이프라인, SaaS 서비스 아키텍처 설계 및 운영 역량 보유\n"
    "• 혁신, 기술 리더십, 비즈니스 임팩트 실현에 중점"
)
summary_p = summary_tf.add_paragraph()
summary_p.text = summary_text
summary_p.font.size = Pt(18)

# -------- 3. 주요 프로젝트 하이라이트 (Projects Grid) --------
slide_projects = prs.slides.add_slide(prs.slide_layouts[5])
slide_projects.shapes.title.text = "Project Highlights"

projects = [
    ("AI Video Generation Framework", 
     "AI 영상 자동 생성 플랫폼 (생산성 60%↑)"),
    ("Edge&Next 의료정보 플랫폼", 
     "실시간 EMR 연동, 고가용성 인증 구조"),
    ("Sales Admin System v3.0", 
     "AI 분석 BI Dashboard, 분석 효율 193%↑"),
    ("MarketLink Survey Center", 
     "KoBERT+GPT LLM 설문 자동 분류, 150배 속도"),
    ("Kakaopay Securities Dashboard", 
     "서비스 트래픽 로그 분석/시각화 자동화"),
    ("교육연수원 e-Learning 시스템", 
     "반응형 웹 UX 설계, 월 사용 30%↑"),
]
table = slide_projects.shapes.add_table(rows=len(projects)+1, cols=2, left=Inches(0.5), top=Inches(1.5), width=Inches(9), height=Inches(3.5)).table
table.columns[0].width = Inches(3.8)
table.columns[1].width = Inches(5.0)
table.cell(0,0).text = "프로젝트명"
table.cell(0,1).text = "핵심 성과/기술"
for i, (prj, note) in enumerate(projects):
    table.cell(i+1,0).text = prj
    table.cell(i+1,1).text = note
for row in range(len(projects)+1):
    for col in range(2):
        p = table.cell(row, col).text_frame.paragraphs[0]
        p.font.size = Pt(14) if row!=0 else Pt(15)
        if row==0:
            p.font.bold = True

# -------- 4. Skills & Technical Stack --------
slide_stack = prs.slides.add_slide(prs.slide_layouts[5])
slide_stack.shapes.title.text = "Technical Skills"

# Use textbox instead of placeholders for reliability
skills_box = slide_stack.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(2))
skills_tf = skills_box.text_frame
skills_tf.word_wrap = True
skills_text = (
    "Backend:  Python (FastAPI), Node.js, PHP7.4\n"
    "Frontend: Vue.js, React, HTML5\n"
    "Database: MySQL, MariaDB, SQLite\n"
    "Infra/DevOps: Nginx, Docker, CentOS, LoadBalancer\n"
    "AI/ML: KoBERT, GPT-4, Claude API, TensorFlow\n"
    "Automation: Google Sheets API, n8n, OpenAPI, BI Tools"
)
skills_p = skills_tf.add_paragraph()
skills_p.text = skills_text
skills_p.font.size = Pt(16)

# -------- 5. 경력 요약 (Career Timeline) --------
slide_career = prs.slides.add_slide(prs.slide_layouts[5])
slide_career.shapes.title.text = "Career Timeline"

career_points = [
    ("2021~현재", "베리타스커넥트 – 생성형 AI 플랫폼/LLM 인프라 구축 (Tech Lead)"),
    ("2018~2021", "아이너스 – 고객소통플랫폼·LMS 시스템 PM/개발"),
    ("2015~2018", "CJ E&M – 미디어 솔루션/글로벌 디지털 전략"),
    ("2009~2015", "메조미디어 – 광고 플랫폼·데이터 분석 시스템 개발"),
    ("2004~2009", "KT-Alpha(구 KT하이텔) IT예산분석관리, Project Management Officer"),
]
table2 = slide_career.shapes.add_table(rows=len(career_points)+1, cols=2, left=Inches(0.5), top=Inches(1.5), width=Inches(9), height=Inches(2.2)).table
table2.columns[0].width = Inches(2)
table2.columns[1].width = Inches(6.8)
table2.cell(0,0).text = "기간"
table2.cell(0,1).text = "주요 경력"
for i, (year, desc) in enumerate(career_points):
    table2.cell(i+1,0).text = year
    table2.cell(i+1,1).text = desc
for row in range(len(career_points)+1):
    for col in range(2):
        p = table2.cell(row, col).text_frame.paragraphs[0]
        p.font.size = Pt(13 if row!=0 else 14)
        if row==0:
            p.font.bold = True

# -------- 6. Performance & Metrics Slide --------
slide_metrics = prs.slides.add_slide(prs.slide_layouts[5])
slide_metrics.shapes.title.text = "Achievements & Metrics"

metrics_box = slide_metrics.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(1.6))
metrics_tf = metrics_box.text_frame
metrics_tf.word_wrap = True
achievements_text = (
    "• 프로젝트 자동화로 평균 작업시간 70% 단축\n"
    "• 데이터 분석 정확도 95% 도달\n"
    "• 반복 업무 자동화율 80% (SaaS형 시스템)\n"
    "• 사용자 수 30%↑, 서버 부하 25%↓ (성공적 시스템 운영 & 최적화)"
)
achievements_p = metrics_tf.add_paragraph()
achievements_p.text = achievements_text
achievements_p.font.size = Pt(18)

# -------- 7. Certifications & Awards --------
slide_cert = prs.slides.add_slide(prs.slide_layouts[5])
slide_cert.shapes.title.text = "Certifications & Awards"
cert_box = slide_cert.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(1.6))
cert_tf = cert_box.text_frame
cert_tf.word_wrap = True
cert_tf.text = (
    "• 정보처리기사 (한국산업인력공단)\n"
    "• 우수 프로젝트상 (메조미디어, 2014)\n"
    "• 기술리더상 (아이너스, 2019)\n"
)

# -------- 8. Vision & Future Goals --------
slide_vision = prs.slides.add_slide(prs.slide_layouts[5])
slide_vision.shapes.title.text = "Vision & Roadmap"
vision_box = slide_vision.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(1.6))
vision_tf = vision_box.text_frame
vision_tf.word_wrap = True
vision_tf.text = (
    "• LLM 기반 데이터 분석 및 시각화 자동화 역량 고도화\n"
    "• 산업 맞춤형 AI 플랫폼 컨설팅\n"
    "• B2B SaaS 시스템 고도화 및 AI Ops 통합 추진"
)

# Save
output_path = "./public/docs/IT_Portfolio_Summary.pptx"
prs.save(output_path)
print(f"PPT 파일이 성공적으로 생성되었습니다: {output_path}")